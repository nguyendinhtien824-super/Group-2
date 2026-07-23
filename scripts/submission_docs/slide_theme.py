from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(14, 31, 53)
NAVY_2 = RGBColor(24, 50, 78)
ORANGE = RGBColor(242, 140, 40)
BLUE = RGBColor(67, 145, 214)
TEAL = RGBColor(42, 157, 143)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(178, 191, 207)
RED = RGBColor(217, 79, 79)


def new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "LAB211 Flash Sale Simulator"
    presentation.core_properties.author = "Nhóm 01"
    return presentation


def add_base_slide(presentation: Presentation, number: int, title: str, eyebrow: str = "LAB211 • NHÓM 01"):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), presentation.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    add_text(slide, eyebrow.upper(), 0.6, 0.32, 8.8, 0.25, 9, ORANGE, bold=True)
    add_text(slide, title, 0.6, 0.65, 11.9, 0.7, 25, WHITE, bold=True)
    add_text(slide, f"{number:02d}", 12.25, 0.38, 0.48, 0.3, 10, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.17), Inches(12.1), Inches(0.02))
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY_2
    footer.line.fill.background()
    return slide


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    size: float,
    color: RGBColor = WHITE,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items: list[str], x: float, y: float, width: float, height: float, size: float = 17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.05)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {item}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(10)
    return box


def add_card(slide, title: str, body: str, x: float, y: float, width: float, height: float, accent: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_2
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.2)
    add_text(slide, title, x + 0.2, y + 0.16, width - 0.4, 0.34, 13, accent, bold=True)
    add_text(slide, body, x + 0.2, y + 0.54, width - 0.4, height - 0.68, 12.5, WHITE)
    return shape


def add_metric(slide, value: str, label: str, x: float, y: float, width: float, accent: RGBColor = ORANGE):
    add_text(slide, value, x, y, width, 0.55, 28, accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.55, width, 0.45, 11, MUTED, align=PP_ALIGN.CENTER)


def add_image_fit(slide, image_path: Path, x: float, y: float, width: float, height: float):
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio > box_ratio:
        target_width = width
        target_height = width / image_ratio
    else:
        target_height = height
        target_width = height * image_ratio
    left = x + (width - target_width) / 2
    top = y + (height - target_height) / 2
    return slide.shapes.add_picture(
        str(image_path), Inches(left), Inches(top), Inches(target_width), Inches(target_height)
    )


def add_table(slide, headers: list[str], rows: list[list[object]], x: float, y: float, width: float, height: float, font_size: float = 11):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(width), Inches(height))
    table = shape.table
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        _style_cell(cell, NAVY, font_size, True)
    for row_index, values in enumerate(rows, start=1):
        for column, value in enumerate(values):
            cell = table.cell(row_index, column)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY_2 if row_index % 2 else RGBColor(31, 61, 91)
            _style_cell(cell, WHITE, font_size, False)
    return shape


def _style_cell(cell, color: RGBColor, size: float, bold: bool) -> None:
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    frame = cell.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
