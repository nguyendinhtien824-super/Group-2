from __future__ import annotations

import xml.etree.ElementTree as element_tree
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from benchmark_data import BenchmarkSummary, create_charts, load_benchmarks, summarize
from slide_theme import (
    BLUE,
    MUTED,
    NAVY,
    NAVY_2,
    ORANGE,
    RED,
    TEAL,
    WHITE,
    add_base_slide,
    add_bullets,
    add_card,
    add_image_fit,
    add_metric,
    add_table,
    add_text,
    new_presentation,
)


def build_slides(project_root: Path, output_path: Path) -> None:
    benchmarks = load_benchmarks(project_root / "data")
    charts = create_charts(benchmarks, project_root / "docs" / "charts")
    summaries = {stock: summarize(rows) for stock, rows in benchmarks.items()}
    test_total, test_failures = _test_count(project_root / "target" / "surefire-reports")
    presentation = new_presentation()

    _cover(presentation)
    _problem(presentation)
    _requirements(presentation)
    _roles(presentation, project_root)
    _architecture(presentation)
    _domain(presentation, project_root)
    _data(presentation, project_root)
    _order(presentation, project_root)
    _race(presentation, project_root)
    _locks(presentation)
    _simulator(presentation, project_root)
    _low_stock(presentation, summaries[100], charts)
    _high_stock(presentation, summaries[2000], charts)
    _verdict(presentation, summaries)
    _quality(presentation, test_total, test_failures)
    _conclusion(presentation)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def _cover(presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.18), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    add_text(slide, "FPT POLYTECHNIC  •  LAB211  •  NHÓM 01", 0.8, 0.7, 11.5, 0.35, 11, ORANGE, bold=True)
    add_text(slide, "FLASH SALE", 0.8, 1.65, 11.6, 0.85, 38, WHITE, bold=True)
    add_text(slide, "CONCURRENCY SIMULATOR", 0.8, 2.45, 11.6, 0.8, 32, BLUE, bold=True)
    add_text(slide, "Race Condition  •  FIFO/Priority  •  4 Lock Mechanisms  •  TPS Benchmark", 0.8, 3.55, 11.5, 0.55, 17, MUTED)
    add_card(slide, "CÂU HỎI NGHIÊN CỨU", "Khóa nào loại bỏ âm kho nhưng vẫn giữ mức giảm throughput không quá 30%?", 0.8, 4.6, 7.2, 1.35, ORANGE)
    add_text(slide, "Thành viên: bổ sung họ tên/MSSV", 8.55, 5.0, 3.7, 0.55, 13, MUTED, align=PP_ALIGN.RIGHT)


def _problem(presentation) -> None:
    slide = add_base_slide(presentation, 2, "Vấn đề: bán vượt kho dưới tải đồng thời")
    add_card(slide, "1. ĐỌC", "Nhiều luồng cùng thấy stock = 1", 0.7, 1.65, 3.7, 1.45, RED)
    add_card(slide, "2. QUYẾT ĐỊNH", "Mọi luồng đều nghĩ có thể bán", 4.8, 1.65, 3.7, 1.45, ORANGE)
    add_card(slide, "3. GHI", "soldQty vượt limited, stock âm", 8.9, 1.65, 3.7, 1.45, RED)
    add_bullets(slide, [
        "Flash Sale dồn hàng nghìn request vào cửa sổ 1–2 giờ.",
        "Đúng nghiệp vụ quan trọng hơn TPS ảo từ dữ liệu sai.",
        "Cần đo cùng workload trên bốn cơ chế và công bố raw run.",
    ], 1.0, 3.8, 11.2, 2.2, 19)


def _requirements(presentation) -> None:
    slide = add_base_slide(presentation, 3, "Phạm vi yêu cầu đã triển khai")
    cards = [
        ("12.527+", "dòng CSV nền", BLUE),
        ("4", "cơ chế đồng bộ", ORANGE),
        ("1000×4×3", "mỗi kịch bản", TEAL),
        ("13", "ngoại lệ nghiệp vụ", RED),
    ]
    for index, (value, label, color) in enumerate(cards):
        add_metric(slide, value, label, 0.75 + index * 3.05, 1.55, 2.55, color)
    add_bullets(slide, [
        "Java 17 • OOP • MVC nhiều lớp • CRUD CSV generic",
        "Event 1–2 giờ • giảm 30–70% • tối đa 2 sản phẩm/customer/product/event",
        "FIFO theo sequence • ưu tiên VIP/PREMIUM nhưng không bỏ qua stock",
        "Biểu đồ TPS • âm kho • retry • kết luận đạt/chưa đạt trung thực",
    ], 0.9, 3.25, 11.6, 2.7, 18)


def _roles(presentation, root: Path) -> None:
    slide = add_base_slide(presentation, 4, "Use Case: bốn vai trò tách quyền")
    add_image_fit(slide, root / "docs" / "use_case_diagram.png", 0.7, 1.45, 7.6, 5.35)
    add_card(slide, "GUEST", "Xem, đăng ký, đăng nhập", 8.65, 1.55, 3.8, 1.0, BLUE)
    add_card(slide, "CUSTOMER", "Mua, voucher, theo dõi/hủy", 8.65, 2.75, 3.8, 1.0, TEAL)
    add_card(slide, "ADMIN", "CRUD, duyệt/hủy, báo cáo", 8.65, 3.95, 3.8, 1.0, ORANGE)
    add_card(slide, "RESEARCHER", "Benchmark và export", 8.65, 5.15, 3.8, 1.0, BLUE)


def _architecture(presentation) -> None:
    slide = add_base_slide(presentation, 5, "Kiến trúc MVC có ranh giới rõ")
    layers = [
        ("VIEW", "Console UI theo vai trò", BLUE),
        ("CONTROLLER", "Điều phối use case", TEAL),
        ("SERVICE", "Invariant • FIFO • lock", ORANGE),
        ("REPOSITORY", "CSV codec • atomic I/O", BLUE),
        ("MODEL", "Entity • enum • version", TEAL),
    ]
    for index, (name, body, color) in enumerate(layers):
        add_card(slide, name, body, 0.85 + index * 2.48, 2.0, 2.15, 2.1, color)
    add_text(slide, "Dependency đi xuống; View không truy cập Repository", 1.2, 5.0, 10.9, 0.5, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)


def _domain(presentation, root: Path) -> None:
    slide = add_base_slide(presentation, 6, "Domain model và versioning")
    add_image_fit(slide, root / "docs" / "class_diagram.png", 0.55, 1.38, 12.2, 5.65)


def _data(presentation, root: Path) -> None:
    slide = add_base_slide(presentation, 7, "Data Generator: đủ lớn, hợp lệ, đọc lại được")
    add_image_fit(slide, root / "docs" / "flowcharts" / "data_generator_flow.png", 0.55, 1.4, 6.6, 5.55)
    add_bullets(slide, [
        "5.000 products • 2.000 customers",
        "10 events • 500 flash items",
        "2.500 orders • 2.500 details • vouchers",
        "Argon2id được CSV quote/escape đúng",
        "Generate → repository round-trip → service boot",
        "Atomic temp + move; UTF-8; đường dẫn tương đối",
    ], 7.4, 1.65, 5.25, 4.8, 16)


def _order(presentation, root: Path) -> None:
    slide = add_base_slide(presentation, 8, "Order flow: kiểm tra lại trong vùng atomic")
    add_image_fit(slide, root / "docs" / "flowcharts" / "order_flow.png", 0.45, 1.35, 7.7, 5.7)
    add_bullets(slide, [
        "Queue sequence + priority tier",
        "Re-read customer, item, event",
        "Fail-fast trạng thái/thời gian",
        "Giới hạn mua nằm trong critical flow",
        "Voucher/wallet/stock rollback khi lỗi",
        "VIP không bao giờ bypass stock",
    ], 8.35, 1.65, 4.25, 4.8, 15.5)


def _race(presentation, root: Path) -> None:
    slide = add_base_slide(presentation, 9, "Race Condition: cùng start, khác chiến lược")
    add_image_fit(slide, root / "docs" / "flowcharts" / "race_condition_flow.png", 0.65, 1.38, 7.25, 5.65)
    add_bullets(slide, [
        "ready/start/done CountDownLatch",
        "No Lock cố ý stale read/write",
        "FileLock dùng sidecar OS lock",
        "synchronized dùng JVM monitor",
        "Optimistic compare version, max 3 retry",
    ], 8.15, 1.85, 4.5, 3.9, 16)


def _locks(presentation) -> None:
    slide = add_base_slide(presentation, 10, "Bốn cơ chế — bốn đánh đổi")
    add_table(slide, ["Cơ chế", "Phạm vi", "Ưu điểm", "Đánh đổi"], [
        ["NO_LOCK", "Không bảo vệ", "Nhanh", "Âm kho/lost update"],
        ["FileLock", "Nhiều process", "An toàn cấp OS", "I/O + lock cost"],
        ["synchronized", "Một JVM", "Đơn giản", "Không cross-process"],
        ["Optimistic", "Version + retry", "Không block", "Cạn retry khi contention"],
    ], 0.75, 1.65, 11.85, 3.55, 13)
    add_card(slide, "TPS", "successCount / elapsedSeconds", 1.0, 5.55, 3.25, 0.95, BLUE)
    add_card(slide, "AN TOÀN", "negativeStock = 0 + consistent", 5.0, 5.55, 3.25, 0.95, TEAL)
    add_card(slide, "MỤC TIÊU", "paired drop ≥ -30%", 9.0, 5.55, 3.25, 0.95, ORANGE)


def _simulator(presentation, root: Path) -> None:
    slide = add_base_slide(presentation, 11, "Simulator và phương pháp thực nghiệm")
    add_image_fit(slide, root / "docs" / "flowcharts" / "simulator_flow.png", 0.5, 1.38, 6.9, 5.7)
    add_card(slide, "A — STOCK 100", "Kích hoạt oversell/race\n1000 × 4 threads × 3", 7.8, 1.75, 4.6, 1.5, RED)
    add_card(slide, "B — STOCK 2000", "Đủ capacity, so TPS\n1000 × 4 threads × 3", 7.8, 3.55, 4.6, 1.5, TEAL)
    add_text(slide, "Raw CSV + Markdown + chart", 7.9, 5.55, 4.4, 0.5, 17, ORANGE, bold=True, align=PP_ALIGN.CENTER)


def _low_stock(presentation, summary: list[BenchmarkSummary], charts: dict[str, Path]) -> None:
    slide = add_base_slide(presentation, 12, "Kết quả A — stock 100: tốc độ ảo vs tính đúng")
    add_image_fit(slide, charts["throughput_100"], 0.55, 1.42, 7.35, 4.6)
    no_lock = _summary(summary, "NO_LOCK")
    add_metric(slide, f"{no_lock.average_tps:.1f}", "No Lock avg TPS", 8.2, 1.75, 1.9, RED)
    add_metric(slide, str(no_lock.negative_stock_total), "tổng âm kho", 10.3, 1.75, 1.9, RED)
    add_card(slide, "SAFE LOCKS", "FileLock / synchronized / optimistic\n0 âm kho • 3/3 consistent", 8.25, 3.15, 4.0, 1.35, TEAL)
    add_card(slide, "VERDICT", "Không cơ chế safe nào đạt ngưỡng TPS ở cả 3 run", 8.25, 4.8, 4.0, 1.2, ORANGE)


def _high_stock(presentation, summary: list[BenchmarkSummary], charts: dict[str, Path]) -> None:
    slide = add_base_slide(presentation, 13, "Kết quả B — stock 2000: capacity công bằng hơn")
    add_image_fit(slide, charts["throughput_2000"], 0.55, 1.42, 7.35, 4.55)
    file_lock = _summary(summary, "FILE_LOCK")
    sync = _summary(summary, "SYNCHRONIZED")
    optimistic = _summary(summary, "OPTIMISTIC")
    add_metric(slide, f"{file_lock.average_tps:.1f}", "FileLock avg TPS", 8.15, 1.55, 1.95, BLUE)
    add_metric(slide, f"{sync.average_tps:.1f}", "sync avg TPS", 10.25, 1.55, 1.95, TEAL)
    add_card(slide, "PAIRED TARGET", f"FileLock: {file_lock.target_runs}/3 run đạt\nsynchronized: {sync.target_runs}/3 run đạt", 8.2, 3.0, 4.0, 1.35, ORANGE)
    add_card(slide, "OPTIMISTIC", f"{optimistic.average_tps:.2f} TPS • {optimistic.retry_total} retries", 8.2, 4.65, 4.0, 1.15, RED)


def _verdict(presentation, summaries: dict[int, list[BenchmarkSummary]]) -> None:
    slide = add_base_slide(presentation, 14, "Kết luận: chưa cơ chế nào đạt xuyên 2 kịch bản")
    high_file = _summary(summaries[2000], "FILE_LOCK")
    high_sync = _summary(summaries[2000], "SYNCHRONIZED")
    add_text(slide, "0", 0.85, 1.45, 3.0, 0.9, 42, RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "safe mechanisms đạt cả 2 kịch bản", 0.65, 2.35, 3.4, 0.5, 13, MUTED, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "Không che giấu raw run và không chọn riêng kết quả đẹp.",
        f"Stock 2000: synchronized {high_sync.target_runs}/3 run đạt; FileLock {high_file.target_runs}/3 run đạt.",
        "Stock 100: cả ba cơ chế an toàn đều 0/3 run đạt ngưỡng TPS.",
        "Optimistic max 3 retry không hợp contention cao trên CSV.",
        "Chọn synchronized cho một JVM; FileLock khi cần phối hợp nhiều process.",
    ], 4.15, 1.45, 8.3, 3.6, 17)
    add_card(slide, "NGHIÊN CỨU TIẾP", "Warm-up • nhiều run • p95 latency • database thật", 4.15, 5.25, 8.0, 1.05, BLUE)


def _quality(presentation, total: int, failures: int) -> None:
    slide = add_base_slide(presentation, 15, "Quality gate, bảo mật và portability")
    add_metric(slide, str(total), "JUnit tests", 0.75, 1.55, 2.4, TEAL)
    add_metric(slide, str(failures), "failure / error", 3.25, 1.55, 2.4, TEAL if failures == 0 else RED)
    add_metric(slide, "Java 17", "fat JAR", 5.75, 1.55, 2.4, BLUE)
    add_metric(slide, "UTF-8", "no BOM", 8.25, 1.55, 2.4, ORANGE)
    add_bullets(slide, [
        "Concurrent last-stock + FIFO/priority + multi-role Customer/Admin",
        "Argon2id; admin secrets qua environment; không hardcode password",
        "Generate → CustomerRepository round-trip → service boot",
        "Đường dẫn tương đối; Maven build; EOF smoke không treo process",
    ], 0.95, 3.45, 11.4, 2.45, 17)


def _conclusion(presentation) -> None:
    slide = add_base_slide(presentation, 16, "Bàn giao: chạy được, đo được, truy vết được")
    add_card(slide, "CODE", "src + test + pom\nMVC • OOP • security", 0.8, 1.55, 3.55, 1.55, BLUE)
    add_card(slide, "EVIDENCE", "CSV raw + Markdown\nchart + 6 diagrams", 4.85, 1.55, 3.55, 1.55, TEAL)
    add_card(slide, "DELIVERABLE", "report.docx + slide.pptx\nREADME + ZIP sạch", 8.9, 1.55, 3.55, 1.55, ORANGE)
    add_text(slide, "DEMO 3 LỆNH", 0.9, 3.75, 3.0, 0.45, 13, ORANGE, bold=True)
    add_bullets(slide, [
        "mvn clean verify",
        "java -jar target/flash-sale-simulator.jar --benchmark-report",
        "java -jar target/flash-sale-simulator.jar",
    ], 0.9, 4.2, 11.4, 1.65, 17)
    add_text(slide, "AI audit log: nhóm ghép bản đã có trước khi nộp.", 0.9, 6.3, 11.4, 0.4, 12, MUTED, align=PP_ALIGN.CENTER)


def _summary(items: list[BenchmarkSummary], lock_type: str) -> BenchmarkSummary:
    return next(item for item in items if item.lock_type == lock_type)


def _test_count(report_dir: Path) -> tuple[int, int]:
    total = 0
    failed = 0
    for report in report_dir.glob("TEST-*.xml"):
        root = element_tree.parse(report).getroot()
        total += int(root.attrib.get("tests", 0))
        failed += int(root.attrib.get("failures", 0)) + int(root.attrib.get("errors", 0))
    if total == 0:
        raise RuntimeError("Chưa có Surefire report; hãy chạy mvn clean verify trước.")
    return total, failed
