import os
import re
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PtInches, Pt as PtSize

def generate_word_report(md_path, docx_path):
    print(f"Starting conversion: {md_path} -> {docx_path}...")
    doc = Document()
    
    # Thiết lập lề trang
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Đọc nội dung Markdown
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    in_table = False
    table_headers = []
    table_rows = []

    for line in lines:
        stripped = line.strip()
        
        # Nhận diện bảng
        if stripped.startswith('|'):
            if '---|' in stripped or '--- |' in stripped:
                # Bỏ qua dòng phân cách bảng
                continue
            
            # Parse các ô của bảng
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            if not in_table:
                in_table = True
                table_headers = cells
            else:
                table_rows.append(cells)
            continue
        else:
            if in_table:
                # Ghi bảng vào tài liệu Word
                table = doc.add_table(rows=1, cols=len(table_headers))
                table.style = 'Light Shading Accent 1'
                
                # Header
                hdr_cells = table.rows[0].cells
                for i, title in enumerate(table_headers):
                    hdr_cells[i].text = title
                    # Format in đậm cho header
                    for paragraph in hdr_cells[i].paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                
                # Rows
                for row_data in table_rows:
                    row_cells = table.add_row().cells
                    for i, val in enumerate(row_data):
                        # Loại bỏ các tag định dạng markdown bold trong ô
                        clean_val = val.replace('**', '')
                        row_cells[i].text = clean_val
                
                doc.add_paragraph() # Dòng trống sau bảng
                in_table = False
                table_headers = []
                table_rows = []

        # Các định dạng dòng thường
        if stripped.startswith('# '):
            title_text = stripped[2:].replace('**', '')
            p = doc.add_heading(level=0)
            run = p.add_run(title_text)
            run.font.name = 'Arial'
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif stripped.startswith('## '):
            heading_text = stripped[3:].replace('**', '')
            p = doc.add_heading(heading_text, level=1)
            p.runs[0].font.name = 'Arial'
        elif stripped.startswith('### '):
            heading_text = stripped[4:].replace('**', '')
            p = doc.add_heading(heading_text, level=2)
            p.runs[0].font.name = 'Arial'
        elif stripped.startswith('#### '):
            heading_text = stripped[5:].replace('**', '')
            p = doc.add_heading(heading_text, level=3)
            p.runs[0].font.name = 'Arial'
        elif stripped.startswith('* ') or stripped.startswith('- '):
            bullet_text = stripped[2:]
            # Xử lý text in đậm trong bullet
            p = doc.add_paragraph(style='List Bullet')
            parts = re.split(r'(\*\*.*?\*\*)', bullet_text)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                else:
                    run = p.add_run(part)
                run.font.name = 'Arial'
        elif stripped == '---' or stripped == '':
            continue
        else:
            p = doc.add_paragraph()
            # Xử lý text in đậm trong đoạn văn
            parts = re.split(r'(\*\*.*?\*\*)', stripped)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                else:
                    run = p.add_run(part)
                run.font.name = 'Arial'

    doc.save(docx_path)
    print(f"Successfully saved {docx_path}!")

def generate_pptx_slide(md_path, pptx_path):
    print(f"Starting conversion: {md_path} -> {pptx_path}...")
    prs = Presentation()
    
    # Đọc nội dung slide markdown
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    slides_content = content.split('\n---\n')
    
    for idx, slide_raw in enumerate(slides_content):
        lines = [line.strip() for line in slide_raw.strip().split('\n') if line.strip()]
        if not lines:
            continue
            
        if idx == 0:
            # Slide tiêu đề chính
            slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
            title_box = slide.shapes.title
            subtitle_box = slide.placeholders[1]
            
            title_text = ""
            sub_text = []
            for line in lines:
                if line.startswith('# '):
                    title_text = line[2:]
                elif line.startswith('## '):
                    sub_text.append(line[3:])
                elif line.startswith('### '):
                    sub_text.append(line[4:])
                else:
                    sub_text.append(line)
            
            title_box.text = title_text
            subtitle_box.text = "\n".join(sub_text)
            
            # Format font
            for p in title_box.text_frame.paragraphs:
                p.font.name = 'Arial'
                p.font.bold = True
            for p in subtitle_box.text_frame.paragraphs:
                p.font.name = 'Arial'
        else:
            # Slide nội dung
            slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
            title_box = slide.shapes.title
            body_box = slide.placeholders[1]
            
            title_text = "Nội dung"
            bullets = []
            
            for line in lines:
                if line.startswith('## '):
                    title_text = line[3:]
                elif line.startswith('* ') or line.startswith('- '):
                    bullets.append(line[2:])
                else:
                    bullets.append(line)
                    
            title_box.text = title_text
            
            # Format font tiêu đề
            for p in title_box.text_frame.paragraphs:
                p.font.name = 'Arial'
                p.font.bold = True
                p.font.size = PtSize(32)
                
            tf = body_box.text_frame
            tf.clear() # Xóa default text
            
            for b in bullets:
                p = tf.add_paragraph()
                p.space_after = PtSize(10)
                
                # Check level/indentation thụt lề
                # Parse markdown bold
                parts = re.split(r'(\*\*.*?\*\*)', b)
                for part in parts:
                    if part.startswith('**') and part.endswith('**'):
                        run = p.add_run()
                        run.text = part[2:-2]
                        run.font.bold = True
                    else:
                        run = p.add_run()
                        run.text = part
                    run.font.name = 'Arial'
                    run.font.size = PtSize(18)
                    
    prs.save(pptx_path)
    print(f"Successfully saved {pptx_path}!")

if __name__ == "__main__":
    current_dir = os.path.dirname(os.path.abspath(__file__))
    
    md_report = os.path.join(current_dir, "report.md")
    docx_report = os.path.join(current_dir, "report.docx")
    
    md_slide = os.path.join(current_dir, "slide.md")
    pptx_slide = os.path.join(current_dir, "slide.pptx")
    
    # Xóa file cũ nếu có trước khi tạo mới
    for path in [docx_report, pptx_slide]:
        if os.path.exists(path):
            os.remove(path)
            print(f"Deleted old file: {path}")
            
    generate_word_report(md_report, docx_report)
    generate_pptx_slide(md_slide, pptx_slide)
